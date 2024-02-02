from datetime import datetime
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import requests
from pptx.util import Inches
import os
import cv2







os.chdir('Code')


doc = Document('../Asset_data.docx')

basic_asset_data = doc.tables[0]
asset_area = basic_asset_data.rows[0].cells[1].text
short_asset_name = basic_asset_data.rows[1].cells[1].text
address = basic_asset_data.rows[2].cells[1].text

asset_description = doc.tables[1].rows[0].cells[0].text
table_rows = []

for row_index, row in enumerate(doc.tables[2].rows):
    table_rows.append([cell.text for cell in row.cells])




airports = []
with open('../Files/airports.csv', 'r', encoding='UTF-8') as f:
    next(f)

    for line in f:
        line = line.replace('\n', '').split(',')
        airports.append(tuple(line))



api_key = ''    # insert api key here
property_name = short_asset_name




def geocode_address(address):
    base_url = 'https://maps.googleapis.com/maps/api/geocode/json'

    params = {
        'address': address,
        'key': api_key
    }

    response = requests.get(base_url, params=params)

    if response.status_code == 200:
        data = response.json()
        if data['status'] == 'OK':
            location = data['results'][0]['geometry']['location']
            lat = location['lat']
            lng = location['lng']
            return lat, lng
        else:
            print('Geocoding failed. Status:', data['status'])
            return None
    else:
        print('HTTP request failed. Status code:', response.status_code)
        return None

# coordinates = geocode_address(address)
# if coordinates:
#     print(f'Coordinates for {address}: Latitude {coordinates[0]}, Longitude {coordinates[1]}')
def find_closest_destination(origin, destinations, api_key):
    base_url = "https://maps.googleapis.com/maps/api/distancematrix/json?"
    params = {
        'destinations': '|'.join(destinations),
        'mode': 'driving',
        'origins': origin,
        'key': api_key
    }
    response = requests.get(base_url, params=params)
    # print(response)
    data = response.json()
    # print(data)

    if 'rows' in data and data['rows']:
        elements = data['rows'][0]['elements']
        distances = [element['distance']['value'] for element in elements if 'distance' in element]
        durations = [element['duration']['value'] for element in elements if 'duration' in element]
        
        if distances:
            min_distance_index = distances.index(min(distances))
            closest_destination = destinations[min_distance_index]
            distance_in_meters = distances[min_distance_index]
            duration_in_seconds = durations[min_distance_index]
            
            hours = duration_in_seconds // 3600
            minutes = (duration_in_seconds % 3600) // 60

            return closest_destination, distance_in_meters, hours, minutes
    return None, None, None, None

origin = address
destinations = [airport for airport, lat, lon in airports]
# print(origin)
# print(destinations)

# closest_destination, distance, hours, minutes
destination1 = list(find_closest_destination(origin, destinations[:20], api_key))
destination2 = list(find_closest_destination(origin, destinations[20:], api_key))
# print(destination1, destination2, sep='\n')

def compare_results(d1, d2):
    t1 = d1[2]*3600 + d1[3]*60
    # print(t1)
    t2 = d2[2]*3600 + d2[3]*60
    # print(t2)

    if t1 <= t2:
        return d1
    else:
        return d2

closest_destination = compare_results(destination1, destination2)
    

if not closest_destination:
    # print(f"The closest destination to {property_name} is {closest_destination[0]} with a distance of {closest_destination[1]} meters.")
    # print(f"Travel time to {closest_destination[0]} is approximately {closest_destination[2]} hours and {closest_destination[3]} minutes.")
# else:
    print("Error retrieving distances and durations.")

locations = [
    [f"{property_name}"] + list(geocode_address(address)),  # Example latitude and longitude for Location A
    [f"{closest_destination[0]}"] + list(geocode_address(closest_destination[0]))  # Example latitude and longitude for Location B
]


def get_static_map_image(locations, size="400x400", maptype="terrain", format="png", scale=1):
    base_url = 'https://maps.googleapis.com/maps/api/staticmap'

    # Calculate the optimal zoom level to fit both locations
    latitudes = [lat for _, lat, _ in locations]
    longitudes = [lon for _, _, lon in locations]
    min_lat = min(latitudes)
    max_lat = max(latitudes)
    min_lon = min(longitudes)
    max_lon = max(longitudes)
    width, height = [int(dim) for dim in size.split('x')]
    # zoom = calculate_zoom_level(min_lat, min_lon, max_lat, max_lon, width, height) - 3
    # print(zoom)

    # Create a list of marker strings for each location
    markers = ['size:mid|color:red|label:{}|{},{}'.format(locations[0][0], locations[0][1], locations[0][2]),
    #  'icon:{}|
     'size:mid|color:red|label:{}|{},{}'.format(#'https://img.icons8.com/?size=50&id=8771&format=png',
            #    'color:blue|label:{}|{},{}'.format(
                   locations[1][0], locations[1][1], locations[1][2])]
                # for label, lat, lng in locations]

    # center = calculate_center(latitudes[0], longitudes[0], latitudes[1], longitudes[1])
    # center = ",".join(map(str, center))

    params = {
        # 'center' : center,
        'markers': markers,
        # 'zoom': zoom,
        'size': size,
        'maptype': maptype,
        'format': format,
        'scale': scale,
        'key': api_key
    }

    response = requests.get(base_url, params=params)

    if response.status_code == 200:
        # Save the image or process it as needed
        with open('../Files/Map.png', 'wb') as f:
            f.write(response.content)
        return '../Files/Map.png'
    else:
        print('HTTP request failed. Status code:', response.status_code)
        return None

size = "800x647"  # Size of the static map image

image_filename = get_static_map_image(locations, size)
# if image_filename:
#     print(f'Static map image saved as {image_filename}')





airport_name, hours, minutes = locations[1][0], closest_destination[2], closest_destination[3]

if int(hours) == 0:
    traveling_hours = ''
elif int(hours) == 1:
    traveling_hours = '1 hour'
elif int(hours) >= 2:
    traveling_hours = f'{hours} hours'

if int(minutes) == 0:
    traveling_minutes = ''
elif int(minutes) == 1:
    if traveling_hours != '':
        traveling_minutes = ' and 1 minute'
    else:
        traveling_minutes = '1 minute'
elif int(minutes) >= 2:
    if traveling_hours != '':
        traveling_minutes = f' and {minutes} minutes'
    else:
        traveling_minutes = f'{minutes} minutes'
    

traveling_time = traveling_hours + traveling_minutes




ppt = Presentation('../Files/Sample_presentation.pptx')


slide = ppt.slides[0]
shape = slide.shapes[10]
text_frame = shape.text_frame
text_frame.clear()  
p = text_frame.add_paragraph()
p.text = asset_area  + ', Italy'
font = p.font
font.name = 'Open Sans Light'
font.size = Pt(20)

paragraph = text_frame.paragraphs[0]
run = paragraph.add_run()
# print(table_rows)
run.text = short_asset_name.upper()
font = run.font
font.bold = True
font.name = 'Open Sans'
font.size = Pt(39)


shape = slide.shapes[11]
text_frame = shape.text_frame
text_frame.clear()  
p = text_frame.paragraphs[0]

today = datetime.now()
current_month = today.strftime("%B")
year = today.year

p.level = 0
p.text = current_month + ' ' + str(year)
font = p.font
font.color.rgb = RGBColor(0, 0, 128)  # Dark blue color
font.bold = True
font.name = 'Open Sans (Body)'
font.size = Pt(14)
# print("Current Month :", current_month)







slides = [(1, 1), (2, 0), (3, 3)]
for t in slides:
    slide = ppt.slides[t[0]]
    shape = slide.shapes[t[1]]

    text_frame = shape.text_frame
    text_frame.clear()  
    p = text_frame.paragraphs[0]
    font = p.font

    p.text = short_asset_name.upper()
    font.color.rgb = RGBColor(12, 158, 217)  # Turquoise, Accent 2
    font.name = 'Open Sans (Body)'
    font.size = Pt(24)




slide = ppt.slides[1]
shape = slide.shapes[5]

# Check if the shape has text
# if shape.has_text_frame:

#     # Select the text frame
text_frame = shape.text_frame

# Clear existing text
text_frame.clear()  

# Add new text
p = text_frame.add_paragraph()
p.text = asset_description + f'\nThe hotel is easily accessible due to the distance of about {traveling_time} from {airport_name}.'
p.line_spacing = 1.5
font = p.font
font.name = 'Open Sans'
font.size = Pt(11)

shape = slide.shapes[7]
text_frame = shape.text_frame
text_frame.clear() 
paragraph = text_frame.paragraphs[0]
run = paragraph.add_run()
run.text = short_asset_name
font = run.font
font.name = 'Open Sans (Body)'
font.size = Pt(10)

shape = slide.shapes[10]
text_frame = shape.text_frame
text_frame.clear() 
paragraph = text_frame.paragraphs[0]
run = paragraph.add_run()
run.text = airport_name
font = run.font
font.name = 'Open Sans (Body)'
font.size = Pt(10)




# Navigate to the slide containing the table. 
# Here, I assume the table is on the first slide (index 0).
slide = ppt.slides[2]

# Find the table on the slide. 
# This code assumes the table is the first shape on the slide.
table = [shape for shape in slide.shapes if shape.has_table][0].table

# Edit a specific cell in the table.
# For instance, to edit the cell at row 1, column 1:
# Access the cell at row 1, column 1.
for i in range(7+1):
    for j in range(2):
        cell = table.cell(i+1, j)
        cell.text = ""

        # Assuming the cell has one paragraph and we want to create a new run in it:
        paragraph = cell.text_frame.paragraphs[0]
        run = paragraph.add_run()
        # print(table_rows)
        run.text = table_rows[i][j]

        # Set the font properties for this run
        font = run.font
        font.name = 'Open Sans'
        font.size = Pt(11)
        font.bold = True if not j else False








def resize_picture(source_path, target_shape_path):
    source = cv2.imread(source_path)
    width, height, _ = source.shape
    target = cv2.imread(target_shape_path)
    target_width, target_height, _ = target.shape

    target_aspect_ratio = target_width / target_height
    new_width = int(height * target_aspect_ratio)

    if new_width <= width:
        x_start = (width - new_width) // 2
        x_end = x_start + new_width
        output = source[x_start:x_end, :]
    else:
        new_height = int(width / target_aspect_ratio)
        y_start = (height - new_height) // 2
        y_end = y_start + new_height
        output = source[:, y_start:y_end]

    cv2.imwrite(source_path[:9] + 'Cropped_images/' + source_path[9:], output)
    # print(source_path[:9] + 'Cropped_images/' + source_path[9:])

def replace_image_in_slide(slide_index, old_image_index, new_image_path, background=False):
    slide = ppt.slides[slide_index]

   
    left = slide.shapes[old_image_index].left
    top = slide.shapes[old_image_index].top
    width = slide.shapes[old_image_index].width
    height = slide.shapes[old_image_index].height
    resize_picture(new_image_path, new_image_path[9:])

    # Delete the old image
    sp = slide.shapes[old_image_index]._element
    sp.getparent().remove(sp)

    # Add the new image at the old image's position
    # print(new_image_path[:9] + 'Cropped_images/' + new_image_path[9:])
    pic = slide.shapes.add_picture(new_image_path[:9] + 'Cropped_images/' + new_image_path[9:], left, top, width=width, height=height)

    if background:
        # print('TRUE!!!')
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    # print(old_image_index)




### SLIDE 1 & 6
new_img_file = f'../Files/Pictures1/Picture.jpg'
replace_image_in_slide(0, 0, new_img_file, True)
replace_image_in_slide(5, 0, new_img_file, True)

### SLIDE 2
new_img_file = f'../Files/Map.png'
replace_image_in_slide(1, 0, new_img_file, True)


### SLIDE 3
new_img_file = f'../Files/Pictures2/Picture.jpg'
replace_image_in_slide(2, 5, new_img_file)

### SLIDE 4
for i in range(7):
    new_img_file = f'../Files/Pictures3/Picture{i + 1}.jpg'
    replace_image_in_slide(3, 4, new_img_file)   # indeces[i]








# Save the modified presentation
ppt.save(f'../{today.year}.{today.month}.{today.day}_COMPANY-NAME_{short_asset_name}.pptx')